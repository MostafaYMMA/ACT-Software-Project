import itertools
import os
import re
import shutil
import sys
import tempfile
import traceback
import zipfile

import win32com.client

# Some real subjects/bodies contain characters (emoji, etc.) outside the
# console's default codepage (cp1252 on Windows) -- without this, the
# [DEBUG] print of such a subject raises UnicodeEncodeError, which the
# per-item try/except in get_approved_cards() catches by skipping that
# email entirely (silent data loss). main.py already does this same
# reconfigure when the app is launched normally; repeated here so this
# module is safe even when run/imported some other way.
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

# ----------------------------------------------------------------------
# Optional third-party extraction libraries (each wrapped so a missing
# library only disables that one format).
# ----------------------------------------------------------------------
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import docx
except ImportError:
    docx = None

try:
    import docx2txt
except ImportError:
    docx2txt = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import extract_msg
except ImportError:
    extract_msg = None

try:
    from odf import text as odf_text
    from odf.opendocument import load as odf_load
except ImportError:
    odf_text = None
    odf_load = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None


# ----------------------------------------------------------------------
# Compiled regular expressions (compiled ONCE, reused everywhere).
# ----------------------------------------------------------------------

# --- Strict approval/timecard logic -- used ONLY for BODY matching
#     (and for informational subject reporting). DO NOT MODIFY. ---
approved_pattern = re.compile(r"\b(?:approved|Submitted|pending|rejected)\b", re.IGNORECASE)

subject_pattern = re.compile(
    r"\btime(?:\s+|-)?card\b|^\s*FW:\s*FYI:?",
    re.IGNORECASE
)

# --- App-to-app sync mail -- NEVER a real timecard/expense email. ---
#
# Mail the other install's app sent to this one (see
# services/outlook_service.py, which owns the full subject grammar and is
# what actually reads these). Matched here only to REJECT: an attachment
# on one of these is a sync payload, and the loose attachment-keyword rule
# in process_email would happily match a current-sheet .xlsx full of
# project/task/hours words and feed the app's own data back through the
# extraction pipeline as if a client had emailed it.
#
# This is a deliberately broad prefix check rather than a copy of
# outlook_service._SUBJECT_PATTERN: anything announcing itself as ACT-SYNC
# should be kept out of here even if it's a version or kind this build
# doesn't recognise.
sync_mail_subject_pattern = re.compile(r"^\s*ACT-SYNC\b", re.IGNORECASE)


def is_sync_mail(item):
    """True if this MailItem is one of the app's own sync messages. Both
    scanners check this first and skip the item entirely."""
    subject = getattr(item, "Subject", "") or ""
    return bool(sync_mail_subject_pattern.match(subject))

# --- Keyword list -- used for ATTACHMENT matching (any ONE keyword
#     is enough) AND for informational reporting in body/subject. ---
KEYWORDS = [
    "Time card",
    "Timecard",
    "Time-card",
    "Time Card Status",
    "Approved",
    "Submitted",
    "Pending",
    "Rejected",
    "Reported time by entry date",
]

KEYWORD_PATTERNS = [
    (kw, re.compile(re.escape(kw), re.IGNORECASE)) for kw in KEYWORDS
]

# --- Expense Report logic -- a SEPARATE filter from the timecard logic
#     above. A match requires ALL THREE terms below to be present
#     (checked across subject + body + attachment text combined), not
#     just one of them like the loose timecard attachment logic. ---
EXPENSE_REQUIRED_TERMS = [
    "Expense Report",
    "Approved",
    "Project",
]

EXPENSE_TERM_PATTERNS = [
    (term, re.compile(re.escape(term), re.IGNORECASE)) for term in EXPENSE_REQUIRED_TERMS
]

# Word-splitting regex (from the uploaded script) -- used to build a
# clean word list from extracted attachment text.
WORD_SPLIT_PATTERN = re.compile(r"\b[\w'-]+\b")

OL_MAIL_ITEM_CLASS = 43
OL_FOLDER_INBOX = 6

DEBUG = True  # set to False to silence per-email debug prints


# ----------------------------------------------------------------------
# Global counters
# ----------------------------------------------------------------------
class Counters:
    def __init__(self):
        self.total_emails = 0
        self.read_emails = 0
        self.unread_emails = 0
        self.subjects_with_approved = 0
        self.subjects_matching_logic = 0          # informational only
        self.total_attachments_scanned = 0
        self.total_matching_attachments = 0        # attachments with >=1 keyword
        self.emails_with_attachment_keywords = 0
        self.emails_with_body_keywords = 0
        self.emails_matched_via_body = 0
        self.emails_matched_via_attachment = 0
        self.emails_with_matching_attachment = 0
        self.total_emails_matched = 0

    def report(self):
        print("\n" + "=" * 70)
        print("SUMMARY")
        print("=" * 70)
        print(f"Total emails scanned:                     {self.total_emails}")
        print(f"Read emails:                               {self.read_emails}")
        print(f"Unread emails:                             {self.unread_emails}")
        print(f"Subjects containing 'Approved' (info only):{self.subjects_with_approved}")
        print(f"Subjects matching logic (info only):       {self.subjects_matching_logic}")
        print(f"Total attachments scanned (all emails):    {self.total_attachments_scanned}")
        print(f"Total matching attachments (>=1 keyword):  {self.total_matching_attachments}")
        print(f"Emails w/ keywords found in attachment:    {self.emails_with_attachment_keywords}")
        print(f"Emails w/ keywords found in body:          {self.emails_with_body_keywords}")
        print(f"Emails matched via BODY (strict logic):    {self.emails_matched_via_body}")
        print(f"Emails matched via ATTACHMENT (keyword):   {self.emails_matched_via_attachment}")
        print(f"Emails with a MATCHING attachment:         {self.emails_with_matching_attachment}")
        print(f"TOTAL EMAILS MATCHED (body OR attachment): {self.total_emails_matched}")
        print("=" * 70)


# ----------------------------------------------------------------------
# Keyword / approval-logic / word-list helpers
# ----------------------------------------------------------------------
def words_from_text(text):
    """
    Split extracted text into a clean list of words (from the uploaded
    reference script). Used as a normalization pass before keyword
    matching on attachment content, so matching is based on the
    actual words present rather than raw substring search alone.
    """
    if not text:
        return []
    return WORD_SPLIT_PATTERN.findall(text)


def find_keywords(text):
    """Return list of configured keywords found in text (case-insensitive)."""
    if not text:
        return []
    found = []
    for kw, pattern in KEYWORD_PATTERNS:
        if pattern.search(text):
            found.append(kw)
    return found


def find_keywords_in_words(word_list):
    """
    Same as find_keywords(), but operates on a pre-split word list by
    rejoining it into a single searchable string first. This lets
    multi-word keywords (e.g. "Time Card Status") still be detected
    even though the source was tokenized, while guaranteeing the
    match is grounded in the actual extracted words.
    """
    if not word_list:
        return []
    rejoined = " ".join(word_list)
    return find_keywords(rejoined)


def matches_approval_logic(text):
    """
    STRICT logic: approved_pattern AND subject_pattern both present.
    Used for BODY matching (and informational subject reporting).
    """
    if not text:
        return False
    return bool(approved_pattern.search(text)) and bool(subject_pattern.search(text))


def find_expense_terms(text):
    """Return list of EXPENSE_REQUIRED_TERMS found in text (case-insensitive)."""
    if not text:
        return []
    found = []
    for term, pattern in EXPENSE_TERM_PATTERNS:
        if pattern.search(text):
            found.append(term)
    return found


def matches_expense_logic(combined_text):
    """
    Expense Report match rule: ALL of EXPENSE_REQUIRED_TERMS
    ("Expense Report", "Approved", "Project") must be present
    somewhere in combined_text. Unlike the timecard attachment
    logic (any ONE keyword), this requires every term.
    """
    found = find_expense_terms(combined_text)
    return len(found) == len(EXPENSE_REQUIRED_TERMS), found


def matches_any_keyword(word_list):
    """
    LOOSE logic: at least ONE keyword from KEYWORDS present in the
    attachment's word list. Used for ATTACHMENT matching.
    """
    return len(find_keywords_in_words(word_list)) > 0


def detect_status(subject, body):
    """
    Which of approved/pending/rejected this email actually is. Subject
    is checked first (e.g. "...Were Approved"), body is the fallback.
    Returns "Approved" / "Submitted" / "Pending" / "Rejected", or None if neither
    text contains one of those words.
    """
    for source in (subject, body):
        match = approved_pattern.search(source or "")
        if match:
            return match.group(0).capitalize()
    return None


# ----------------------------------------------------------------------
# Per-format text extraction functions
# ----------------------------------------------------------------------
def extract_pdf_text(filepath):
    if pdfplumber is None:
        return ""
    text_chunks = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_chunks.append(page_text)
    except Exception:
        traceback.print_exc()
    return "\n".join(text_chunks)


def extract_docx_text(filepath):
    text_chunks = []
    try:
        if docx is not None:
            document = docx.Document(filepath)
            for para in document.paragraphs:
                if para.text:
                    text_chunks.append(para.text)
            for table in document.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text_chunks.append(cell.text)
            if text_chunks:
                return "\n".join(text_chunks)
    except Exception:
        traceback.print_exc()

    try:
        if docx2txt is not None:
            return docx2txt.process(filepath) or ""
    except Exception:
        traceback.print_exc()

    return "\n".join(text_chunks)


def extract_doc_text(filepath):
    try:
        if docx2txt is not None:
            result = docx2txt.process(filepath)
            if result:
                return result
    except Exception:
        pass

    try:
        with open(filepath, "rb") as f:
            raw = f.read()
        decoded = raw.decode("latin-1", errors="ignore")
        return re.sub(r"[^\x20-\x7E\n\r\t]+", " ", decoded)
    except Exception:
        traceback.print_exc()
        return ""


def extract_excel_text(filepath, ext):
    text_chunks = []
    try:
        if ext == ".xlsx" and openpyxl is not None:
            wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            text_chunks.append(str(cell))
            wb.close()
        elif ext == ".xls" and xlrd is not None:
            wb = xlrd.open_workbook(filepath)
            for sheet in wb.sheets():
                for row_idx in range(sheet.nrows):
                    for cell in sheet.row(row_idx):
                        if cell.value not in (None, ""):
                            text_chunks.append(str(cell.value))
    except Exception:
        traceback.print_exc()
    return "\n".join(text_chunks)


def extract_powerpoint_text(filepath, ext):
    text_chunks = []
    try:
        if ext == ".pptx" and Presentation is not None:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_chunks.append(shape.text)
        elif ext == ".ppt":
            with open(filepath, "rb") as f:
                raw = f.read()
            decoded = raw.decode("latin-1", errors="ignore")
            text_chunks.append(re.sub(r"[^\x20-\x7E\n\r\t]+", " ", decoded))
    except Exception:
        traceback.print_exc()
    return "\n".join(text_chunks)


def extract_msg_text(filepath):
    if extract_msg is None:
        return ""
    text_chunks = []
    try:
        msg = extract_msg.Message(filepath)
        if msg.subject:
            text_chunks.append(msg.subject)
        if msg.body:
            text_chunks.append(msg.body)
        msg.close()
    except Exception:
        traceback.print_exc()
    return "\n".join(text_chunks)


def extract_odf_text(filepath):
    if odf_load is None or odf_text is None:
        return ""
    text_chunks = []
    try:
        doc = odf_load(filepath)
        for element in doc.getElementsByType(odf_text.P):
            text_chunks.append(str(element))
    except Exception:
        traceback.print_exc()
    return "\n".join(text_chunks)


def extract_html_text(filepath):
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            raw_html = f.read()
        if BeautifulSoup is not None:
            soup = BeautifulSoup(raw_html, "lxml")
            return soup.get_text(separator="\n")
        return re.sub(r"<[^>]+>", " ", raw_html)
    except Exception:
        traceback.print_exc()
        return ""


def extract_xml_text(filepath):
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            raw_xml = f.read()
        if BeautifulSoup is not None:
            soup = BeautifulSoup(raw_xml, "xml")
            return soup.get_text(separator="\n")
        return re.sub(r"<[^>]+>", " ", raw_xml)
    except Exception:
        traceback.print_exc()
        return ""


def extract_rtf_text(filepath):
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            raw_rtf = f.read()
        if rtf_to_text is not None:
            return rtf_to_text(raw_rtf)
        return raw_rtf
    except Exception:
        traceback.print_exc()
        return ""


def extract_plain_text(filepath):
    try:
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        traceback.print_exc()
        return ""


def extract_csv_text(filepath):
    """From the uploaded script: dedicated CSV extraction via csv module."""
    import csv
    text_chunks = []
    try:
        with open(filepath, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                text_chunks.append(" ".join(row))
    except Exception:
        traceback.print_exc()
    return "\n".join(text_chunks)


def extract_zip_text(filepath):
    text_chunks = []
    temp_extract_dir = tempfile.mkdtemp(prefix="zip_extract_")
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            for member in zf.namelist():
                try:
                    extracted_path = zf.extract(member, temp_extract_dir)
                    if os.path.isdir(extracted_path):
                        continue
                    member_text = extract_text_from_file(extracted_path)
                    if member_text:
                        text_chunks.append(member_text)
                except Exception:
                    traceback.print_exc()
    except Exception:
        traceback.print_exc()
    finally:
        shutil.rmtree(temp_extract_dir, ignore_errors=True)
    return "\n".join(text_chunks)


def extract_text_from_file(filepath):
    """Detect file type by extension and dispatch to the right extractor."""
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == ".pdf":
            return extract_pdf_text(filepath)
        elif ext == ".docx":
            return extract_docx_text(filepath)
        elif ext == ".doc":
            return extract_doc_text(filepath)
        elif ext == ".xlsx":
            return extract_excel_text(filepath, ext)
        elif ext == ".xls":
            return extract_excel_text(filepath, ext)
        elif ext == ".pptx":
            return extract_powerpoint_text(filepath, ext)
        elif ext == ".ppt":
            return extract_powerpoint_text(filepath, ext)
        elif ext == ".csv":
            return extract_csv_text(filepath)
        elif ext in (".txt", ".log"):
            return extract_plain_text(filepath)
        elif ext in (".html", ".htm"):
            return extract_html_text(filepath)
        elif ext == ".xml":
            return extract_xml_text(filepath)
        elif ext == ".rtf":
            return extract_rtf_text(filepath)
        elif ext == ".msg":
            return extract_msg_text(filepath)
        elif ext == ".zip":
            return extract_zip_text(filepath)
        elif ext in (".odt", ".ods", ".odp"):
            return extract_odf_text(filepath)
        else:
            return ""
    except Exception:
        traceback.print_exc()
        return ""


# ----------------------------------------------------------------------
# Attachment processing
# ----------------------------------------------------------------------
def process_attachment(attachment, temp_dir, counters):
    """
    Save one attachment to disk, extract its text, split it into a
    clean word list (words_from_text), and check whether that word
    list contains AT LEAST ONE of the KEYWORDS (loose match -- this
    is what decides an attachment-based email match). Returns a
    result dict; deletes the temp file before returning.
    """
    result = {
        "filename": None,
        "filetype": None,
        "keywords_found": [],
        "text": "",
        "word_count": 0,
        "matches_keyword": False,   # True if ANY keyword found
    }

    filename = None
    saved_path = None
    try:
        filename = attachment.FileName
        ext = os.path.splitext(filename)[1].lower()
        safe_name = f"{abs(hash(filename))}_{filename}"
        saved_path = os.path.join(temp_dir, safe_name)

        attachment.SaveAsFile(saved_path)
        counters.total_attachments_scanned += 1

        text = extract_text_from_file(saved_path)

        # Normalize via word-list split (from the uploaded reference script)
        word_list = words_from_text(text)
        keywords_found = find_keywords_in_words(word_list)

        result["filename"] = filename
        result["filetype"] = ext if ext else "(unknown)"
        result["keywords_found"] = keywords_found
        result["text"] = text
        result["word_count"] = len(word_list)
        result["matches_keyword"] = len(keywords_found) > 0

        if result["matches_keyword"]:
            counters.total_matching_attachments += 1

        if DEBUG:
            print(f"        [DEBUG] Attachment '{filename}' -> "
                  f"words={result['word_count']} | "
                  f"keywords={keywords_found} | "
                  f"matches_keyword={result['matches_keyword']}")

    except Exception:
        traceback.print_exc()
    finally:
        if saved_path and os.path.exists(saved_path):
            try:
                os.remove(saved_path)
            except Exception:
                traceback.print_exc()

    return result


def extract_attachment_text_only(attachment, temp_dir):
    """
    Same save-to-disk / extract / cleanup flow as process_attachment(),
    but returns just (filename, extracted_text) with no dependency on
    the timecard Counters object or KEYWORDS list. Used by the
    Expense Report scan so it stays independent of the timecard logic.
    """
    filename = None
    saved_path = None
    text = ""
    try:
        filename = attachment.FileName
        safe_name = f"{abs(hash(filename))}_{filename}"
        saved_path = os.path.join(temp_dir, safe_name)
        attachment.SaveAsFile(saved_path)
        text = extract_text_from_file(saved_path)
    except Exception:
        traceback.print_exc()
    finally:
        if saved_path and os.path.exists(saved_path):
            try:
                os.remove(saved_path)
            except Exception:
                traceback.print_exc()
    return filename, text


def process_email_expense(item, temp_dir):
    """
    Process a single MailItem for the Expense Report filter.

    MATCH RULE (all three required, checked against subject + body +
    every attachment's extracted text combined):
        "Expense Report", "Approved", "Project"

    Returns a dict describing the matched email (subject, sender,
    received, where each required term was found) if all three terms
    are present somewhere in the email, or None if not matched.
    """
    # Same early rejection as process_email -- see is_sync_mail.
    if is_sync_mail(item):
        return None

    subject = getattr(item, "Subject", "") or ""

    try:
        sender = item.SenderName
    except Exception:
        sender = "(unknown sender)"

    try:
        received = item.ReceivedTime
    except Exception:
        received = "(unknown date)"

    try:
        body = item.Body or ""
    except Exception:
        body = ""

    attachment_texts = []
    try:
        attachments = item.Attachments
        for i in range(1, attachments.Count + 1):
            att_filename, att_text = extract_attachment_text_only(attachments.Item(i), temp_dir)
            if att_text:
                attachment_texts.append((att_filename, att_text))
    except Exception:
        traceback.print_exc()

    combined_text = "\n".join([subject, body] + [t for _, t in attachment_texts])

    is_matched, terms_found = matches_expense_logic(combined_text)

    if DEBUG:
        print(f"[DEBUG][Expense] Subject: '{subject}' | terms_found={terms_found} | "
              f"matched={is_matched}")

    if not is_matched:
        return None

    # Where each required term showed up, for reporting.
    term_locations = []
    for term, pattern in EXPENSE_TERM_PATTERNS:
        found_in = []
        if pattern.search(subject):
            found_in.append("Subject")
        if pattern.search(body):
            found_in.append("Body")
        for att_filename, att_text in attachment_texts:
            if pattern.search(att_text):
                found_in.append(f"Attachment '{att_filename}'")
        term_locations.append((term, found_in))

    print(f"\n>>> EXPENSE REPORT MATCH: {subject}")
    print("-" * 70)
    print(f"Subject:   {subject}")
    print(f"Sender:    {sender}")
    print(f"Received:  {received}")
    print("Required terms found in:")
    for term, found_in in term_locations:
        print(f"    - {term}: {', '.join(found_in) if found_in else '(not found here)'}")
    print("-" * 70)

    return {
        "subject": subject,
        "sender": sender,
        "received": received,
        "terms_found": terms_found,
        "term_locations": term_locations,
        "mail_item": item,
    }


# ----------------------------------------------------------------------
# Email processing
# ----------------------------------------------------------------------
def process_email(item, temp_dir, counters):
    """
    Process a single MailItem. Every attachment is opened and scanned
    regardless of subject content.

    MATCH RULE:
      - Subject alone NEVER triggers a match (informational only).
      - BODY triggers a match only via the STRICT logic
        (approved_pattern AND subject_pattern both present).
      - ATTACHMENT triggers a match if its word list contains AT
        LEAST ONE keyword from KEYWORDS (loose match).

    Returns:
        A dict describing the matched email (with all the detail
        needed for reporting) if this email matched, or None if it
        did not match. This return value is what get_approved_cards()
        collects into its matching_emails list.
    """
    # Before anything else, and before a single attachment is opened: the
    # app's own sync mail is handled by outlook_service.scan_sync_mails /
    # sync_service.pull_updates, and must never be run through approval
    # detection or extraction here.
    if is_sync_mail(item):
        if DEBUG:
            print(f"[DEBUG] Skipping app sync mail: '{getattr(item, 'Subject', '')}'")
        return None

    counters.total_emails += 1

    try:
        if item.UnRead:
            counters.unread_emails += 1
        else:
            counters.read_emails += 1
    except Exception:
        traceback.print_exc()

    subject = getattr(item, "Subject", "") or ""

    # Subject checks -- INFO/DEBUG/counters only, never affects match.
    has_approved = bool(approved_pattern.search(subject))
    if has_approved:
        counters.subjects_with_approved += 1

    has_subject_pattern = bool(subject_pattern.search(subject))
    subject_would_have_matched = has_approved and has_subject_pattern
    if subject_would_have_matched:
        counters.subjects_matching_logic += 1

    if DEBUG:
        print(f"[DEBUG] Subject: '{subject}' | approved={has_approved} | "
              f"pattern={has_subject_pattern} | "
              f"(subject alone does NOT trigger match)")

    try:
        sender = item.SenderName
    except Exception:
        sender = "(unknown sender)"

    try:
        received = item.ReceivedTime
    except Exception:
        received = "(unknown date)"

    try:
        body = item.Body or ""
    except Exception:
        body = ""

    body_keywords = find_keywords(body)
    found_in_body = len(body_keywords) > 0
    if found_in_body:
        counters.emails_with_body_keywords += 1

    # --- BODY match uses the STRICT approval logic ---
    body_approval_match = matches_approval_logic(body)

    if DEBUG:
        print(f"        [DEBUG] Body strict approval_logic_match={body_approval_match}")

    subject_keywords = find_keywords(subject)

    # --- ALWAYS open every attachment ---
    attachment_results = []
    found_in_attachment = False
    attachment_keyword_match = False
    try:
        attachments = item.Attachments
        if DEBUG:
            print(f"        [DEBUG] Attachment count: {attachments.Count}")
        for i in range(1, attachments.Count + 1):
            attachment = attachments.Item(i)
            att_result = process_attachment(attachment, temp_dir, counters)
            if att_result["keywords_found"]:
                found_in_attachment = True
            if att_result["matches_keyword"]:
                attachment_keyword_match = True
            attachment_results.append(att_result)
    except Exception:
        traceback.print_exc()

    if found_in_attachment:
        counters.emails_with_attachment_keywords += 1

    if attachment_keyword_match:
        counters.emails_with_matching_attachment += 1

    # --- FINAL MATCH DECISION: body (strict) OR attachment (loose keyword) ---
    is_matched = body_approval_match or attachment_keyword_match
    if not is_matched:
        return None  # not a match -- skip reporting, nothing to collect

    counters.total_emails_matched += 1
    if body_approval_match:
        counters.emails_matched_via_body += 1
    if attachment_keyword_match:
        counters.emails_matched_via_attachment += 1

    status = detect_status(subject, body)
    if status is None:
        # Some matched emails (e.g. forwarded invoices) carry the actual
        # Approved/Pending/Rejected word only inside a matching
        # attachment's own text, not in the email's own subject/body --
        # detect_status() can't see that on its own. Without this
        # fallback, extract() would tag every entry from that email
        # status=None and save_cards() would silently drop all of them.
        for att in attachment_results:
            if att.get("keywords_found"):
                status = detect_status("", att.get("text") or "")
                if status:
                    break

    matched_via = []
    if body_approval_match:
        matched_via.append("BODY (strict logic)")
    if attachment_keyword_match:
        matched_via.append("ATTACHMENT (keyword match)")

    # --- Print the matched subject clearly ---
    print(f"\n>>> MATCHED EMAIL SUBJECT: {subject}")

    print("-" * 70)
    print(f"Subject:          {subject}")
    print(f"Sender:           {sender}")
    print(f"Received:         {received}")
    print(f"Matched via:      {' + '.join(matched_via)}")
    print(f"(Subject alone matched approval logic: {subject_would_have_matched} "
          f"-- informational only, did not cause this match)")

    locations = []
    if subject_keywords:
        locations.append(f"Subject ({', '.join(subject_keywords)})")
    if found_in_body:
        locations.append(f"Body ({', '.join(body_keywords)})")
    for att in attachment_results:
        if att["keywords_found"]:
            locations.append(
                f"Attachment '{att['filename']}' [{att['filetype']}] "
                f"({', '.join(att['keywords_found'])})"
            )

    if locations:
        print("Keywords found in:")
        for loc in locations:
            print(f"    - {loc}")
    else:
        print("Keywords found in: (none)")

    # --- Attachments processed summary ---
    if attachment_results:
        print("Attachments processed:")
        for att in attachment_results:
            flag = " <-- MATCHED (keyword)" if att["matches_keyword"] else ""
            print(f"    - {att['filename']} ({att['filetype']}, {att['word_count']} words){flag}")
    else:
        print("Attachments processed: (none)")

    # --- Full breakdown: every attachment name + every keyword found in it ---
    print("Attachment name -> keywords found (detailed):")
    if attachment_results:
        for att in attachment_results:
            name = att["filename"] if att["filename"] else "(unknown filename)"
            if att["keywords_found"]:
                kw_list = ", ".join(att["keywords_found"])
            else:
                kw_list = "(no keywords found)"
            print(f"    - {name}: {kw_list}")
    else:
        print("    (no attachments on this email)")

    print("-" * 70)

    # Build and return the record for this matched email.
    matched_email = {
        "subject": subject,
        "sender": sender,
        "received": received,
        "status": status,                    # "Approved" / "Pending" / "Rejected" / None
        "matched_via": matched_via,          # e.g. ["BODY (strict logic)"]
        "matched_via_body": body_approval_match,
        "matched_via_attachment": attachment_keyword_match,
        "subject_keywords": subject_keywords,
        "body_keywords": body_keywords,
        "keyword_locations": locations,
        "attachments": [
            {
                "filename": att["filename"],
                "filetype": att["filetype"],
                "word_count": att["word_count"],
                "keywords_found": att["keywords_found"],
                "matches_keyword": att["matches_keyword"],
                "text": att["text"],
            }
            for att in attachment_results
        ],
        "mail_item": item,  # raw Outlook MailItem, in case caller needs it
    }

    return matched_email


# ----------------------------------------------------------------------
# Folder connection helper
# ----------------------------------------------------------------------
def get_outlook_folder(namespace, folder_name="Inbox"):
    """
    Connect helper adapted from the uploaded reference script: returns
    the Inbox by default, or looks up a named subfolder under Inbox.
    """
    inbox = namespace.GetDefaultFolder(OL_FOLDER_INBOX)
    if folder_name.lower() == "inbox":
        return inbox
    for folder in inbox.Folders:
        if folder.Name.lower() == folder_name.lower():
            return folder
    raise ValueError(f"Folder '{folder_name}' not found under Inbox.")


# ----------------------------------------------------------------------
# Public entry point: get_approved_cards()
# ----------------------------------------------------------------------
def _received_time_naive(item):
    """
    Best-effort naive datetime.datetime for item.ReceivedTime, used to
    compare against start_date/end_date ahead of the (much slower)
    process_email() call. Outlook/pywin32 hands back ReceivedTime as a
    tz-aware pywintypes.datetime whose tzinfo doesn't reflect an actual
    conversion (it's local wall-clock time labeled as a fixed offset) --
    dropping tzinfo here matches how storage_service._parse_received()
    already treats the same value once it's stored as the "received"
    column, so the two stay consistent.
    """
    try:
        received = item.ReceivedTime
    except Exception:
        return None
    return received.replace(tzinfo=None) if getattr(received, "tzinfo", None) is not None else received


def get_newest_received(folder_name="Inbox"):
    """
    ReceivedTime of the newest item currently in the folder, as a naive
    datetime, or None if the folder is empty / Outlook can't be reached.

    Read BEFORE a scan starts, not after, and stored as the new high-water
    mark once that scan finishes (see storage_service.set_last_scan_time
    and sync_service.sync_cards). Taking it up front is what makes mail
    that lands DURING a long scan get picked up by the next one instead of
    being stepped over: the mark never claims to cover anything the scan
    couldn't have seen.

    Only the first few items are inspected -- the collection is sorted
    newest-first, so the loop is really just skipping the odd item that has
    no ReceivedTime at all (some non-mail items), not searching.
    """
    try:
        outlook = win32com.client.Dispatch("Outlook.Application")
        namespace = outlook.GetNamespace("MAPI")
        folder = get_outlook_folder(namespace, folder_name)
    except Exception as exc:
        print(f"Could not read the newest received time: {exc}")
        return None

    items = folder.Items
    items.Sort("[ReceivedTime]", True)  # newest first
    for item in itertools.islice(items, 10):
        received_at = _received_time_naive(item)
        if received_at is not None:
            return received_at
    return None


def get_approved_cards(folder_name="Inbox", print_report=True, limit=None,
                        start_date=None, end_date=None):
    """
    Connect to Outlook, scan the given folder (Inbox by default) using
    the same single-pass logic as main(), and return the list of
    matched emails.

    A match is an email where either:
      - the BODY satisfies the strict approval logic (approved_pattern
        AND subject_pattern both present), or
      - at least one ATTACHMENT contains one of the KEYWORDS.

    Args:
        folder_name: name of the folder under Inbox to scan
                     ("Inbox" itself, or a named subfolder).
        print_report: if True, prints the Counters summary at the end
                     (same output as main()).
        limit: if set, only scan the `limit` most recently received
                     emails instead of the whole folder. Ignored for
                     items outside [start_date, end_date] when those are
                     given -- see below.
        start_date/end_date: if given, only emails whose ReceivedTime
                     falls within [start_date, end_date] (inclusive) are
                     scanned/matched. Since the folder is sorted
                     newest-first, an email older than start_date ends
                     the scan immediately rather than counting against
                     `limit`.

    Returns:
        matching_emails: list of dicts, one per matched email, in the
        same format built inside process_email() (subject, sender,
        received, matched_via, keyword_locations, attachments, etc).
    """
    counters = Counters()
    matching_emails = []

    try:
        outlook = win32com.client.Dispatch("Outlook.Application")
        namespace = outlook.GetNamespace("MAPI")
        folder = get_outlook_folder(namespace, folder_name)
    except Exception as exc:
        print(f"FATAL: Could not connect to Outlook: {exc}")
        traceback.print_exc()
        return matching_emails

    items = folder.Items
    items.Sort("[ReceivedTime]", True)  # newest first, so limit takes the most recent
    temp_dir = tempfile.mkdtemp(prefix="outlook_scan_")

    try:
        # Single pass over the folder (capped at `limit` if given) -- O(n).
        for item in (itertools.islice(items, limit) if limit else items):
            try:
                if getattr(item, "Class", None) != OL_MAIL_ITEM_CLASS:
                    continue

                if start_date is not None or end_date is not None:
                    received_at = _received_time_naive(item)
                    if received_at is None:
                        continue
                    if end_date is not None and received_at > end_date:
                        continue  # newer than the window -- keep scanning
                    if start_date is not None and received_at < start_date:
                        break  # newest-first order: nothing after this is in range

                matched_email = process_email(item, temp_dir, counters)
                if matched_email is not None:
                    matching_emails.append(matched_email)
            except Exception:
                traceback.print_exc()
                continue
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    if print_report:
        counters.report()

    return matching_emails


def get_expense_reports(folder_name="Inbox", print_report=True, limit=None,
                         start_date=None, end_date=None):
    """
    Connect to Outlook, scan the given folder (Inbox by default), and
    return the list of emails matching the Expense Report filter.

    A match is an email where "Expense Report", "Approved", AND
    "Project" are ALL present somewhere across the subject, body, and
    attachment text combined.

    Args mirror get_approved_cards() (folder_name, print_report, limit,
    start_date/end_date).

    Returns:
        matching_emails: list of dicts, one per matched email
        (subject, sender, received, terms_found, term_locations).
    """
    matching_emails = []
    total_scanned = 0

    try:
        outlook = win32com.client.Dispatch("Outlook.Application")
        namespace = outlook.GetNamespace("MAPI")
        folder = get_outlook_folder(namespace, folder_name)
    except Exception as exc:
        print(f"FATAL: Could not connect to Outlook: {exc}")
        traceback.print_exc()
        return matching_emails

    items = folder.Items
    items.Sort("[ReceivedTime]", True)  # newest first
    temp_dir = tempfile.mkdtemp(prefix="outlook_expense_scan_")

    try:
        for item in (itertools.islice(items, limit) if limit else items):
            try:
                if getattr(item, "Class", None) != OL_MAIL_ITEM_CLASS:
                    continue

                if start_date is not None or end_date is not None:
                    received_at = _received_time_naive(item)
                    if received_at is None:
                        continue
                    if end_date is not None and received_at > end_date:
                        continue
                    if start_date is not None and received_at < start_date:
                        break

                total_scanned += 1
                matched_email = process_email_expense(item, temp_dir)
                if matched_email is not None:
                    matching_emails.append(matched_email)
            except Exception:
                traceback.print_exc()
                continue
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    if print_report:
        print("\n" + "=" * 70)
        print("EXPENSE REPORT SUMMARY")
        print("=" * 70)
        print(f"Total emails scanned:            {total_scanned}")
        print(f"Expense Reports found:           {len(matching_emails)}")
        if matching_emails:
            print("Matching subjects:")
            for m in matching_emails:
                print(f"    - {m['subject']}")
        print("=" * 70)

    return matching_emails


# ----------------------------------------------------------------------
# Main entry point
# ----------------------------------------------------------------------
def main():
    matching_emails = get_approved_cards("Inbox")
    print(f"\n{len(matching_emails)} matching email(s) returned by get_approved_cards().")

    expense_emails = get_expense_reports("Inbox")
    print(f"\n{len(expense_emails)} Expense Report email(s) returned by get_expense_reports().")


if __name__ == "__main__":
    main()